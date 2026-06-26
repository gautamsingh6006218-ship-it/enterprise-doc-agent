"""
tests/ingestion/test_ingestion_service.py

What does this cover?
- LoaderRegistry: routing, registration, supported extensions.
- IngestionService: all formats (txt, md, html, docx, pptx, xlsx, pdf).
- RBAC fields: owner_id, access_roles, visibility, tenant_id.
- Error handling: missing file, unsupported format.
- Deduplication: file_hash consistency, unique IDs.
- SmartPdfLoader: text detection and routing logic.
- UnstructuredLoader: DOCX, PPTX, XLSX extraction via Unstructured.
"""

from pathlib import Path

import pytest

from agent.ingestion.loader_registry import LoaderRegistry
from agent.ingestion.loaders.base import BaseDocumentLoader
from agent.ingestion.models import Document
from agent.services.ingestion_service import IngestionService


# ── Fixtures ───────────────────────────────────────────────────────────────────

@pytest.fixture()
def service() -> IngestionService:
    return IngestionService()


@pytest.fixture()
def sample_txt(tmp_path: Path) -> Path:
    p = tmp_path / "notes.txt"
    p.write_text("Enterprise document content.\nSecond line.", encoding="utf-8")
    return p


@pytest.fixture()
def sample_md(tmp_path: Path) -> Path:
    p = tmp_path / "readme.md"
    p.write_text("# Title\n\nSome markdown content.", encoding="utf-8")
    return p


@pytest.fixture()
def sample_html(tmp_path: Path) -> Path:
    p = tmp_path / "page.html"
    p.write_text(
        "<html><head><title>My Page</title></head>"
        "<body><p>Hello HTML</p><script>alert(1)</script></body></html>",
        encoding="utf-8",
    )
    return p


@pytest.fixture()
def sample_docx(tmp_path: Path) -> Path:
    from docx import Document as DocxDocument
    doc = DocxDocument()
    doc.add_paragraph("Enterprise DOCX content.")
    doc.add_paragraph("Second paragraph.")
    path = tmp_path / "report.docx"
    doc.save(str(path))
    return path


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide content for enterprise RAG."
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def sample_xlsx(tmp_path: Path) -> Path:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Department", "Score"])
    ws.append(["Alice", "Engineering", 95])
    ws.append(["Bob", "HR", 88])
    path = tmp_path / "data.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture()
def sample_pdf_text(tmp_path: Path) -> Path:
    """Minimal valid single-page PDF with extractable text."""
    content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Contents 4 0 R/Resources<</Font<</F1<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>>>>>>>endobj
4 0 obj<</Length 44>>
stream
BT /F1 12 Tf 100 700 Td (Hello PDF) Tj ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000274 00000 n
trailer<</Size 5/Root 1 0 R>>
startxref
368
%%EOF"""
    p = tmp_path / "sample.pdf"
    p.write_bytes(content)
    return p


# ── LoaderRegistry ─────────────────────────────────────────────────────────────

class TestLoaderRegistry:

    def test_all_expected_extensions_registered(self):
        registry = LoaderRegistry()
        expected = [
            ".pdf", ".txt", ".md", ".markdown",
            ".html", ".htm", ".docx", ".pptx",
            ".xlsx", ".csv", ".png", ".jpg", ".eml",
        ]
        for ext in expected:
            assert ext in registry.supported_extensions, f"Missing: {ext}"

    def test_pdf_routes_to_smart_loader(self):
        from agent.ingestion.loaders.pdf_smart import SmartPdfLoader
        loader = LoaderRegistry().get_loader("doc.pdf")
        assert isinstance(loader, SmartPdfLoader)

    def test_docx_routes_to_unstructured(self):
        from agent.ingestion.loaders.unstructured_loader import UnstructuredLoader
        loader = LoaderRegistry().get_loader("doc.docx")
        assert isinstance(loader, UnstructuredLoader)

    def test_pptx_routes_to_unstructured(self):
        from agent.ingestion.loaders.unstructured_loader import UnstructuredLoader
        loader = LoaderRegistry().get_loader("deck.pptx")
        assert isinstance(loader, UnstructuredLoader)

    def test_txt_routes_to_txt_loader(self):
        from agent.ingestion.loaders.txt import TxtLoader
        loader = LoaderRegistry().get_loader("notes.txt")
        assert isinstance(loader, TxtLoader)

    def test_html_routes_to_html_loader(self):
        from agent.ingestion.loaders.html import HtmlLoader
        loader = LoaderRegistry().get_loader("page.html")
        assert isinstance(loader, HtmlLoader)

    def test_unsupported_extension_raises(self):
        with pytest.raises(ValueError, match="No loader registered for '.xyz'"):
            LoaderRegistry().get_loader("file.xyz")

    def test_register_overrides_existing(self):
        class FakeLoader(BaseDocumentLoader):
            @property
            def supported_extensions(self) -> list[str]:
                return [".pdf"]
            def load(self, file_path: str, tenant_id: str = "default") -> Document:
                raise NotImplementedError

        registry = LoaderRegistry()
        registry.register(FakeLoader())
        assert isinstance(registry.get_loader("doc.pdf"), FakeLoader)


# ── Format loading ─────────────────────────────────────────────────────────────

class TestFormats:

    def test_txt(self, service: IngestionService, sample_txt: Path):
        doc = service.ingest(str(sample_txt)).document
        assert doc.source_type == "txt"
        assert "Enterprise document content" in doc.text
        assert doc.title == "notes"

    def test_markdown(self, service: IngestionService, sample_md: Path):
        doc = service.ingest(str(sample_md)).document
        assert doc.source_type == "markdown"
        assert "markdown content" in doc.text

    def test_html_strips_scripts(self, service: IngestionService, sample_html: Path):
        doc = service.ingest(str(sample_html)).document
        assert "alert" not in doc.text
        assert "Hello HTML" in doc.text
        assert doc.title == "My Page"

    def test_docx(self, service: IngestionService, sample_docx: Path):
        result = service.ingest(str(sample_docx))
        assert result.success is True
        assert "Enterprise DOCX content" in result.document.text
        assert result.document.source_type == "docx"

    def test_pptx(self, service: IngestionService, sample_pptx: Path):
        result = service.ingest(str(sample_pptx))
        assert result.success is True
        doc = result.document
        assert doc.source_type == "pptx"
        assert len(doc.text) > 0

    def test_xlsx(self, service: IngestionService, sample_xlsx: Path):
        result = service.ingest(str(sample_xlsx))
        assert result.success is True
        doc = result.document
        assert doc.source_type == "xlsx"
        assert len(doc.text) > 0

    def test_pdf_text_based(self, service: IngestionService, sample_pdf_text: Path):
        result = service.ingest(str(sample_pdf_text))
        assert result.success is True
        assert result.document.source_type == "pdf"


# ── RBAC ───────────────────────────────────────────────────────────────────────

class TestRBAC:

    def test_default_rbac_fields(self, service: IngestionService, sample_txt: Path):
        doc = service.ingest(str(sample_txt)).document
        assert doc.owner_id == "system"
        assert doc.access_roles == []
        assert doc.visibility == "public"
        assert doc.tenant_id == "default"

    def test_custom_owner(self, service: IngestionService, sample_txt: Path):
        doc = service.ingest(str(sample_txt), owner_id="user-42").document
        assert doc.owner_id == "user-42"

    def test_access_roles(self, service: IngestionService, sample_txt: Path):
        doc = service.ingest(str(sample_txt), access_roles=["hr", "legal"]).document
        assert set(doc.access_roles) == {"hr", "legal"}

    def test_visibility_restricted(self, service: IngestionService, sample_txt: Path):
        doc = service.ingest(str(sample_txt), visibility="restricted").document
        assert doc.visibility == "restricted"

    def test_tenant_isolation(self, service: IngestionService, sample_txt: Path):
        a = service.ingest(str(sample_txt), tenant_id="tenant-a").document
        b = service.ingest(str(sample_txt), tenant_id="tenant-b").document
        assert a.tenant_id != b.tenant_id


# ── Error handling ─────────────────────────────────────────────────────────────

class TestErrors:

    def test_missing_file(self, service: IngestionService):
        result = service.ingest("/non/existent/file.txt")
        assert result.success is False
        assert "not found" in result.error.lower()

    def test_unsupported_format(self, service: IngestionService, tmp_path: Path):
        f = tmp_path / "data.xyz"
        f.write_text("data")
        result = service.ingest(str(f))
        assert result.success is False
        assert ".xyz" in result.error


# ── Deduplication ──────────────────────────────────────────────────────────────

class TestDeduplication:

    def test_same_file_same_hash(self, service: IngestionService, sample_txt: Path):
        h1 = service.ingest(str(sample_txt)).document.file_hash
        h2 = service.ingest(str(sample_txt)).document.file_hash
        assert h1 == h2

    def test_unique_id_per_ingest(self, service: IngestionService, sample_txt: Path):
        id1 = service.ingest(str(sample_txt)).document.id
        id2 = service.ingest(str(sample_txt)).document.id
        assert id1 != id2

    def test_different_content_different_hash(
        self, service: IngestionService, tmp_path: Path
    ):
        f1 = tmp_path / "a.txt"
        f2 = tmp_path / "b.txt"
        f1.write_text("content A")
        f2.write_text("content B")
        assert (
            service.ingest(str(f1)).document.file_hash
            != service.ingest(str(f2)).document.file_hash
        )


# ── SmartPdfLoader routing ─────────────────────────────────────────────────────

class TestSmartPdfLoader:

    def test_text_pdf_returns_pymupdf_result(self, tmp_path: Path):
        """
        When text density check passes, SmartPdfLoader returns the PdfLoader result
        without falling through to Docling or OCR.
        """
        from unittest.mock import patch
        from agent.ingestion.loaders.pdf_smart import SmartPdfLoader

        path = tmp_path / "sample.pdf"
        content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Contents 4 0 R/Resources<</Font<</F1<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>>>>>>>endobj
4 0 obj<</Length 44>>
stream
BT /F1 12 Tf 100 700 Td (Hello PDF) Tj ET
endstream
endobj
xref
0 5
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000274 00000 n
trailer<</Size 5/Root 1 0 R>>
startxref
368
%%EOF"""
        path.write_bytes(content)

        loader = SmartPdfLoader()
        # Mock _is_text_sparse to return False — simulates a text-rich PDF.
        # This isolates the routing logic from PDF rendering behaviour.
        with patch.object(loader, "_is_text_sparse", return_value=False):
            result = loader.load(str(path))

        assert result.source_type == "pdf"
        assert result.metadata.get("extraction_method") == "pymupdf4llm_markdown"

    def test_sparse_pdf_triggers_fallback(self, tmp_path: Path):
        """Sparse text PDF should add extraction_warning when no advanced loaders available."""
        from agent.ingestion.loaders.pdf_smart import SmartPdfLoader

        loader = SmartPdfLoader()
        loader._docling_loader = None
        loader._ocr_loader = None

        # Use the minimal hand-crafted PDF (produces very little text from pymupdf4llm)
        content = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R>>endobj
xref
0 4
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
trailer<</Size 4/Root 1 0 R>>
startxref
190
%%EOF"""
        path = tmp_path / "sparse.pdf"
        path.write_bytes(content)

        result = loader.load(str(path))
        assert result.source_type == "pdf"
        assert "extraction_warning" in result.metadata
